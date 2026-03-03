from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

src = r"C:/Users/Administrator/.openclaw/workspace/SISTRA/docs/Guia_SISTRA_PC_Horizontal_Paso_a_Paso.pptx"
out = r"C:/Users/Administrator/.openclaw/workspace/SISTRA/docs/Guia_SISTRA_PC_Horizontal_Paso_a_Paso_PRO.pptx"

prs = Presentation(src)
SW, SH = prs.slide_width, prs.slide_height

GREEN_DARK = RGBColor(0x00,0x5C,0x3C)
GREEN_SOFT = RGBColor(0xE8,0xF4,0xEF)
GREEN_PALE = RGBColor(0xF4,0xFA,0xF7)
ACCENT = RGBColor(0x0B,0x8F,0x63)
TEXT_DARK = RGBColor(0x12,0x2B,0x22)
WHITE = RGBColor(0xFF,0xFF,0xFF)

step_labels = {
    2: "PASO 1 · LOGIN",
    3: "PASO 2 · TABLERO",
    4: "PASO 3 · TRÁMITES",
    5: "PASO 4 · NUEVA CAPTURA",
    6: "PASO 5 · DETALLE",
    7: "PASO 6 · CONFIGURACIÓN",
    8: "CIERRE · BUENAS PRÁCTICAS"
}

for idx, slide in enumerate(prs.slides, start=1):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GREEN_PALE

    if idx == 1:
        bg.fore_color.rgb = GREEN_DARK
        if len(slide.shapes) >= 2:
            title = slide.shapes[0]
            subtitle = slide.shapes[1]
            title.left = int(SW*0.07); title.top = int(SH*0.28); title.width = int(SW*0.55); title.height = int(SH*0.22)
            subtitle.left = int(SW*0.07); subtitle.top = int(SH*0.52); subtitle.width = int(SW*0.55); subtitle.height = int(SH*0.14)

            for p in title.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = 'Aptos Display'
                    r.font.size = Pt(42)
                    r.font.bold = True
                    r.font.color.rgb = WHITE
            for p in subtitle.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = 'Aptos'
                    r.font.size = Pt(18)
                    r.font.color.rgb = RGBColor(0xD9,0xF0,0xE7)

        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(SW*0.07), int(SH*0.14), int(SW*0.26), int(SH*0.08))
        chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT
        chip.line.fill.background()
        chip.text_frame.text = "Guía de Capacitación Institucional"
        for p in chip.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name='Aptos'; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=WHITE

        deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(SW*0.70), int(SH*0.16), int(SW*0.22), int(SW*0.22))
        deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0x1A,0x8D,0x66)
        deco.line.fill.background()
        deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(SW*0.78), int(SH*0.28), int(SW*0.14), int(SW*0.14))
        deco2.fill.solid(); deco2.fill.fore_color.rgb = RGBColor(0x46,0xAF,0x8D)
        deco2.line.fill.background()

    else:
        title = slide.shapes[0]
        title.left = int(SW*0.03); title.top = int(SH*0.02); title.width = int(SW*0.64); title.height = int(SH*0.09)
        title.fill.solid(); title.fill.fore_color.rgb = GREEN_DARK
        title.line.fill.background()
        for p in title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = 'Aptos Display'
                r.font.size = Pt(26)
                r.font.bold = True
                r.font.color.rgb = WHITE

        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(SW*0.70), int(SH*0.02), int(SW*0.25), int(SH*0.06))
        tag.fill.solid(); tag.fill.fore_color.rgb = ACCENT
        tag.line.fill.background()
        tag.text_frame.text = step_labels.get(idx, "PASO")
        for p in tag.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name='Aptos'; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=WHITE

        if len(slide.shapes) > 1:
            pic = slide.shapes[1]
            pic.left = int(SW*0.04); pic.top = int(SH*0.17); pic.width = int(SW*0.63); pic.height = int(SH*0.72)
            try:
                pic.line.color.rgb = RGBColor(0xB7,0xD8,0xCB)
                pic.line.width = Pt(1.2)
            except Exception:
                pass

        target = slide.shapes[2] if len(slide.shapes) > 2 else slide.shapes[1]
        target.left = int(SW*0.70); target.top = int(SH*0.17); target.width = int(SW*0.27); target.height = int(SH*0.72)
        target.fill.solid(); target.fill.fore_color.rgb = GREEN_SOFT
        target.line.color.rgb = RGBColor(0x9D,0xC9,0xB8)
        target.line.width = Pt(1.2)

        tf = target.text_frame
        tf.margin_left = Pt(12); tf.margin_right = Pt(12); tf.margin_top = Pt(10); tf.margin_bottom = Pt(10)
        original = tf.text
        if "Qué se edita" not in original:
            tf.clear()
            p0 = tf.paragraphs[0]; p0.text = "Qué se edita:"; p0.level = 0
            p1 = tf.add_paragraph(); p1.text = "Parámetros o datos del trámite según pantalla"; p1.level = 1
            p2 = tf.add_paragraph(); p2.text = "Qué se captura:"; p2.level = 0
            p3 = tf.add_paragraph(); p3.text = "Campos obligatorios del solicitante y validaciones"; p3.level = 1
            p4 = tf.add_paragraph(); p4.text = "Indicaciones del paso:"; p4.level = 0
            p5 = tf.add_paragraph(); p5.text = original.replace('•','').replace('¿Qué hacer aquí?','').strip(); p5.level = 1

        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = 'Aptos'
                r.font.size = Pt(13 if p.level==0 else 11)
                r.font.bold = (p.level==0)
                r.font.color.rgb = TEXT_DARK

s8 = prs.slides[7]
if len(s8.shapes)>=2:
    body = s8.shapes[1]
    body.left = int(SW*0.08); body.top = int(SH*0.20); body.width = int(SW*0.84); body.height = int(SH*0.64)
    body.fill.solid(); body.fill.fore_color.rgb = GREEN_SOFT
    body.line.color.rgb = RGBColor(0x9D,0xC9,0xB8)
    tf = body.text_frame
    tf.margin_left=Pt(20); tf.margin_right=Pt(20); tf.margin_top=Pt(16)
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name='Aptos'; r.font.size=Pt(18 if p.level==0 else 15); r.font.color.rgb=TEXT_DARK

prs.save(out)
chk = Presentation(out)
print(f'OK slides={len(chk.slides)} file={out}')
